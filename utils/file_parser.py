# utils/file_parser.py
"""
Universal File Parser
Handles 25+ file formats with pluggable parser functions.
Goal: Reliable, production-safe, easily extendable.
"""

import os
import json
import csv
import zipfile
import tarfile
import gzip
import sqlite3
import yaml
import pandas as pd
import pyarrow.parquet as pq
import pyarrow.feather as feather
import avro.datafile
import avro.io
import docx
import pptx
import openpyxl
import xlrd
import pdfplumber
import email
import xml.etree.ElementTree as ET
from bs4 import BeautifulSoup
from pydub.utils import mediainfo

SUPPORTED_FORMATS = [
    ".csv", ".tsv", ".xlsx", ".xls", ".json", ".yaml", ".yml", ".xml",
    ".html", ".txt", ".log", ".sql", ".parquet", ".feather", ".avro",
    ".pdf", ".eml", ".hl7", ".docx", ".doc", ".rtf", ".odt", ".ods",
    ".pptx", ".ppt", ".gz", ".zip", ".tar", ".sqlite", ".db", ".wav",
    ".mp3", ".mp4"
]

class FileParser:
    def __init__(self):
        self.parsers = {
            ".csv": self.parse_csv,
            ".tsv": lambda f: self.parse_csv(f, delimiter="\t"),
            ".xlsx": self.parse_excel,
            ".xls": self.parse_excel,
            ".json": self.parse_json,
            ".yaml": self.parse_yaml,
            ".yml": self.parse_yaml,
            ".xml": self.parse_xml,
            ".html": self.parse_html,
            ".txt": self.parse_text,
            ".log": self.parse_text,
            ".sql": self.parse_text,
            ".parquet": self.parse_parquet,
            ".feather": self.parse_feather,
            ".avro": self.parse_avro,
            ".pdf": self.parse_pdf,
            ".eml": self.parse_eml,
            ".hl7": self.parse_text,
            ".docx": self.parse_docx,
            ".doc": self.parse_docx,
            ".rtf": self.parse_text,  # can be replaced with an RTF parser
            ".odt": self.parse_text,  # lightweight ODT parse
            ".ods": self.parse_excel,
            ".pptx": self.parse_pptx,
            ".ppt": self.parse_pptx,
            ".gz": self.parse_gzip,
            ".zip": self.parse_zip,
            ".tar": self.parse_tar,
            ".sqlite": self.parse_sqlite,
            ".db": self.parse_sqlite,
            ".wav": self.parse_audio_metadata,
            ".mp3": self.parse_audio_metadata,
            ".mp4": self.parse_video_metadata
        }

    def parse(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()
        if ext not in self.parsers:
            raise ValueError(f"Unsupported file format: {ext}")
        return self.parsers[ext](file_path)

    # --- Text-based ---
    def parse_text(self, file_path, encoding="utf-8"):
        with open(file_path, "r", encoding=encoding, errors="ignore") as f:
            return f.read()

    def parse_csv(self, file_path, delimiter=","):
        return pd.read_csv(file_path, delimiter=delimiter).to_dict(orient="records")

    def parse_excel(self, file_path):
        return pd.read_excel(file_path).to_dict(orient="records")

    def parse_json(self, file_path):
        with open(file_path, "r", encoding="utf-8") as f:
            return json.load(f)

    def parse_yaml(self, file_path):
        with open(file_path, "r", encoding="utf-8") as f:
            return yaml.safe_load(f)

    def parse_xml(self, file_path):
        tree = ET.parse(file_path)
        return {tree.getroot().tag: self._xml_to_dict(tree.getroot())}

    def _xml_to_dict(self, elem):
        return {
            elem.tag: {
                "attributes": elem.attrib,
                "text": elem.text.strip() if elem.text else "",
                "children": [self._xml_to_dict(e) for e in elem]
            }
        }

    def parse_html(self, file_path):
        with open(file_path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f, "html.parser")
        return soup.get_text()

    # --- Structured Data ---
    def parse_parquet(self, file_path):
        return pq.read_table(file_path).to_pandas().to_dict(orient="records")

    def parse_feather(self, file_path):
        return feather.read_feather(file_path).to_dict(orient="records")

    def parse_avro(self, file_path):
        with open(file_path, "rb") as f:
            reader = avro.datafile.DataFileReader(f, avro.io.DatumReader())
            return [record for record in reader]

    # --- Documents ---
    def parse_pdf(self, file_path):
        text = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() + "\n"
        return text.strip()

    def parse_eml(self, file_path):
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            msg = email.message_from_file(f)
        return {
            "subject": msg.get("subject"),
            "from": msg.get("from"),
            "to": msg.get("to"),
            "body": msg.get_payload()
        }

    def parse_docx(self, file_path):
        doc = docx.Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])

    def parse_pptx(self, file_path):
        prs = pptx.Presentation(file_path)
        slides_text = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            slides_text.append("\n".join(slide_text))
        return slides_text

    # --- Archives ---
    def parse_gzip(self, file_path):
        with gzip.open(file_path, "rt", encoding="utf-8", errors="ignore") as f:
            return f.read()

    def parse_zip(self, file_path):
        data = {}
        with zipfile.ZipFile(file_path, "r") as zip_ref:
            for name in zip_ref.namelist():
                with zip_ref.open(name) as f:
                    try:
                        data[name] = f.read().decode("utf-8", errors="ignore")
                    except:
                        data[name] = f.read()
        return data

    def parse_tar(self, file_path):
        data = {}
        with tarfile.open(file_path, "r") as tar_ref:
            for member in tar_ref.getmembers():
                f = tar_ref.extractfile(member)
                if f:
                    try:
                        data[member.name] = f.read().decode("utf-8", errors="ignore")
                    except:
                        data[member.name] = f.read()
        return data

    # --- DB ---
    def parse_sqlite(self, file_path):
        conn = sqlite3.connect(file_path)
        cursor = conn.cursor()
        tables = cursor.execute("SELECT name FROM sqlite_master WHERE type='table'").fetchall()
        db_data = {}
        for (table_name,) in tables:
            rows = cursor.execute(f"SELECT * FROM {table_name}").fetchall()
            db_data[table_name] = rows
        conn.close()
        return db_data

    # --- Media ---
    def parse_audio_metadata(self, file_path):
        return mediainfo(file_path)

    def parse_video_metadata(self, file_path):
        return mediainfo(file_path)
        
